# -*- coding: utf-8 -*-
"""
Created on Wed Jun 29 14:45:33 2022

@author: Aluneth
"""

import json
import os
import nltk, re, pprint
import operator
import math
import time
import math

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from nltk import word_tokenize
from nltk.text import TextCollection
from tqdm import tqdm #进度条库
from math import log

dir = r"F:\Document_Understanding_Project_with_NLP_technique\111"
file_name = 'result.jsonl' #jsonl - jsonlist
outputfile = '234.jsonl'
final_result_file = 'FinalResult.jsonl'
final_anayzle_result = 'FinalAnalyzeResult.jsonl'
sum_noun_file = "SumNounFile.jsonl"
stop_word_removed = "StopWordRemoved.jsonl"
stop_word_removed_2 = "StopWordRemoved2.jsonl"

def translate_ppt_to_json(filename):

    prs = Presentation(filename)
    txt = ""
    
    for x in range(len(prs.slides)):
        for shape in prs.slides[x].shapes:
            if hasattr(shape, "text"):
                row_text = shape.text.encode('utf-8').strip().decode()
                txt += row_text
                txt += ", "
        group_shapes = [shp for shp in prs.slides[x].shapes
                        if shp.shape_type == MSO_SHAPE_TYPE.GROUP]
        for group_shape in group_shapes:
            for shape in group_shape.shapes:
                if shape.has_text_frame:
                    row_text = shape.text.encode('utf-8').strip().decode()
                    txt += row_text
                    
                    txt += ", "
        txt += ". "
    return txt.replace("\n", "").replace("\r", "")

def extract_keyword_from_txt(file_name) :
    #抽取关键词， 仅限名词
    with open(file_name, 'r', encoding='utf-8') as file_obj, open(outputfile, 'w', encoding ='utf-8') as file_o :
        for line in file_obj :
            _json = json.loads(line)
            text = _json['Text']
            path = _json['Filename']
            result = txt(text)  
            noun_save = {"FilePath": path,
                     "Noun_Result": result
                     }
            
            file_o.write(json.dumps(noun_save, ensure_ascii = False))
            file_o.write("\n")
            
def get_the_stop_word(filename) :
    #得到stopword
    
    with open (filename,'r', encoding = 'utf-8') as file_obj, open(sum_noun_file,'w', encoding = 'utf-8') as file_o :
        sum_dic = {}
       # del_dict = {}
        for line in file_obj :
            _json = json.loads(line)
            
            temp_dict = _json['Counter Result']
            for key in temp_dict.keys() :
                if key in sum_dic :
                    sum_dic[key] += 1
                else :
                    sum_dic[key] = 1
        sorted_dict = dict(sorted(sum_dic.items(), key = operator.itemgetter(1), reverse = True)) 
        file_o.write(json.dumps(sorted_dict,ensure_ascii = False))        
   
def remove_stop_words(filename) :
    #移除stop word
    word_list = []
    temp_dic = {}
    temp_loc_dict = {}
    
    idf_number_100 = 2.5702664987391572

    with open (filename,'r', encoding = 'utf-8') as file_obj, open(stop_word_removed,'w', encoding = 'utf-8') as file_o : 
        for line in file_obj :
            _json = json.loads(line)
            temp_loc_dict = _json['Counter Result']
            for key in temp_loc_dict :
                word_list.append(key)
        corpus = TextCollection(word_list)
        #print(len(word_list))
            
        for i in tqdm(word_list) :
            
            idf = corpus.idf(i) #求这个词在整个文件里出现的次数
            temp_dic[i] = idf
            
        for key in tqdm(list(temp_dic.keys())) :
            if temp_dic.get(key) <= idf_number_100 :
                temp_dic.pop(key)

                
        result_dic = {
                      "Counter Result":temp_dic
                }
        file_o.write(json.dumps(result_dic,ensure_ascii = False)) 
        file_o.write("\n")

def remove_stop_words_2(filename1, filename2):
    with open (filename1,'r',encoding='utf-8') as file_obj, open(filename2,'r',encoding='utf-8') as file_ob, open(stop_word_removed_2,'w', encoding = 'utf-8') as file_o :
        word_dic = json.load(file_ob)
        for key in list(word_dic.keys()) :  #不能递归移除,把key转为列表存储再移除
            if word_dic.get(key) <= 100 :
                word_dic.pop(key)
        for line in file_obj :
            _json = json.loads(line)
            filepath = _json['FileName']
            temp_dic = _json['Counter Result']
            for key in list(temp_dic.keys()) :
                if key in word_dic :
                    temp_dic.pop(key)
            result_dic = {"FileName":filepath,
                          "Counter Result":temp_dic
                }
            file_o.write(json.dumps(result_dic,ensure_ascii = False)) 
            file_o.write("\n")
        
    

   
def countthe_number_of_nouns_in_the_file(filename):
    with open(filename,'r',encoding='utf-8') as file_obj, open(final_result_file,'w',encoding = 'utf-8') as file_o :
        for line in file_obj :
            _json = json.loads(line)
            noun_text = _json['Noun_Result']
            pathh = _json['FilePath']
            analyze_result = noun_counter(noun_text)
            save_result = {"FileName": pathh,
                           "Counter Result": analyze_result
                           }
            file_o.write(json.dumps(save_result, ensure_ascii= False))
            file_o.write("\n")
            
def get_the_top_K_word(filename) :
    #找出现K次的名词
    K = 23
    noun_list = []
    with open(filename,'r',encoding='utf-8' ) as file_obj :
        for line in file_obj :
            _json = json.loads(line)
            word_dic = _json['Counter Result']
            for key in word_dic :
                word = key
                if word_dic[key] == K :
                    noun_list.append(word)
    print(noun_list)
                                    
     
def noun_counter(txt_txt):
    counter = {}
    sum1 = 1
    for key in txt_txt :
        word = key #提取单词
        if (len(word) > 1):   #移除长度为1的标点符号
            if key in counter :  
                i = counter.get(word) #遍历比对重复数量，如果有则+1
                i += 1
                counter[word] = i
            else:
                counter.update({word:sum1}) #如果比对完成就保存去下一个单词
    
    sorted_dict = dict(sorted(counter.items(), key = operator.itemgetter(1), reverse = True)) #value从大到小排序，itemgetter（1）是比对value，（0）是比对key

    return sorted_dict


def remove_key(dict1, keyword):
    dict1 = dict1.pop(keyword)
    return dict1

     
    
def txt(text) :
     sentences = nltk.sent_tokenize(text)
     sentences = [nltk.word_tokenize(sent) for sent in sentences]
     sentences = [nltk.pos_tag(sent) for sent in sentences]#将内容句子分块为单词并标注单词类型
     Noun_sentence = []
     grammar = "NPCHUNK:{<NN|NNP|NNS|NNPS>?<NN|NNP|NNS|NNPS>}"   #名词后面只要跟名词就视作一个块，直到不再跟名词为止
     cp = nltk.RegexpParser(grammar)
     for s in sentences :
         tree = cp.parse(s)
         for subtree in tree.subtrees():
             if subtree.label() == 'NPCHUNK' : 
                 np = ' '.join(w for w, t in subtree.leaves())  #因为分块出的块实际上是树，所以要把叶子抽出来重新写进list中
                 Noun_sentence.append(np)
     return Noun_sentence

def analyze_file_and_remove_stop_word(dict1) :
    result_dict = {}
    word_list = []
    for key in dict1 :
        word_list.append(key)
    corpus = TextCollection(word_list)
    #print(corpus)
    idf = corpus.idf('be')
    print(idf)


def extract_ppt(filepath):
    with open(file_name, 'w', encoding = 'utf-8' ) as file_obj: 
        for file in os.listdir(dir) :
            filepath = os.path.join(dir, file)
            txt = translate_ppt_to_json(filepath) #执行抽取文件的方法
            _json = {"Filename": filepath, 
                     "Text": txt, 
                     }#json文件分为txt和filepath两部分
            file_obj.write(json.dumps(_json, ensure_ascii = False))
            file_obj.write("\n")

def cacluate_log(some):
    number = log(1320/(some+1))
    print(number)
        
#extract_keyword_from_txt(file_name)   
#countthe_number_of_nouns_in_the_file(outputfile)
#get_the_top_K_word(final_anayzle_result)
#get_the_stop_word(final_anayzle_result) 
#remove_stop_words_2(final_anayzle_result, sum_noun_file)
#remove_stop_words(final_anayzle_result)
cacluate_log(20)